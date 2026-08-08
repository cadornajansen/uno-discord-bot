from dataclasses import dataclass
import logging
from pathlib import Path
from typing import Any, Optional

import pdf_inspector
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class DocumentAnalysis:
    """Holds extracted Markdown text, metadata, and warnings from a document."""

    filename: str
    file_type: str
    markdown: str
    warnings: tuple[str, ...]
    page_count: Optional[int] = None
    slide_count: Optional[int] = None


class DocumentError(Exception):
    """Base exception for document processing operations."""

    pass


class UnsupportedFileError(DocumentError):
    """Raised when an unsupported file extension is provided."""

    pass


class DocumentParseError(DocumentError):
    """Raised when PDF or PPTX file parsing fails due to corruption or encryption."""

    pass


def format_markdown_table(table: Any) -> str:
    """Convert a python-pptx Table shape into a Markdown table string.

    Args:
        table: pptx.table.Table object.

    Returns:
        Formatted Markdown table string.
    """
    rows = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            # Replace internal newlines with space to preserve table structure
            cell_text = cell.text.replace("\n", " ").strip() if cell.text else ""
            row_cells.append(cell_text)
        rows.append(row_cells)

    if not rows:
        return ""

    num_cols = len(rows[0])
    md_lines = []

    # Header row
    md_lines.append("| " + " | ".join(rows[0]) + " |")
    # Separator row
    md_lines.append("| " + " | ".join(["---"] * num_cols) + " |")

    # Data rows
    for r in rows[1:]:
        # Pad row cells if row is shorter than header
        padded = r + [""] * (num_cols - len(r))
        md_lines.append("| " + " | ".join(padded[:num_cols]) + " |")

    return "\n".join(md_lines)


class DocumentService:
    """Service handling text extraction and normalization for PDF and PPTX files."""

    def __init__(self, max_chars: int = 50000):
        self.max_chars = max_chars

    async def extract(self, path: Path, filename: str) -> DocumentAnalysis:
        """Extract and normalize textual content from a local PDF or PPTX file.

        Args:
            path: Absolute Path to the local document file.
            filename: Original filename for labeling and validation.

        Returns:
            DocumentAnalysis dataclass instance.

        Raises:
            UnsupportedFileError: If extension is not .pdf or .pptx.
            DocumentParseError: If file is corrupt or unreadable.
        """
        ext = path.suffix.lower()

        if ext == ".pdf":
            analysis = self._extract_pdf(path, filename)
        elif ext == ".pptx":
            analysis = self._extract_pptx(path, filename)
        else:
            raise UnsupportedFileError(
                f"Unsupported file format '{ext}'. Uno currently supports only .pdf and .pptx files."
            )

        # Enforce maximum character budget
        if len(analysis.markdown) > self.max_chars:
            truncated_md = analysis.markdown[: self.max_chars]
            new_warnings = list(analysis.warnings) + [
                f"This document exceeds Uno's current analysis limit. Only the first {self.max_chars:,} characters were analyzed."
            ]
            return DocumentAnalysis(
                filename=analysis.filename,
                file_type=analysis.file_type,
                markdown=truncated_md,
                warnings=tuple(new_warnings),
                page_count=analysis.page_count,
                slide_count=analysis.slide_count,
            )

        return analysis

    def _extract_pdf(self, path: Path, filename: str) -> DocumentAnalysis:
        """Private helper for PDF extraction using pdf-inspector."""
        try:
            result = pdf_inspector.process_pdf(str(path))
        except Exception as e:
            logger.error(f"pdf-inspector failed to process '{filename}': {e}")
            raise DocumentParseError(
                f"Failed to process PDF '{filename}'. It may be corrupted or encrypted."
            ) from e

        raw_md = getattr(result, "markdown", "") or ""
        pages_needing_ocr = getattr(result, "pages_needing_ocr", []) or []

        warnings = []
        if pages_needing_ocr:
            ocr_count = len(pages_needing_ocr)
            warnings.append(
                f"{ocr_count} page(s) contain scanned or image-based content requiring OCR and were not fully readable."
            )

        if not raw_md.strip():
            warnings.append(
                "This PDF appears to contain mostly scanned or image-based pages. OCR support is not available yet."
            )

        return DocumentAnalysis(
            filename=filename,
            file_type="PDF",
            markdown=raw_md,
            warnings=tuple(warnings),
        )

    def _extract_pptx(self, path: Path, filename: str) -> DocumentAnalysis:
        """Private helper for PPTX extraction using python-pptx."""
        try:
            prs = Presentation(str(path))
        except Exception as e:
            logger.error(f"python-pptx failed to open '{filename}': {e}")
            raise DocumentParseError(
                f"Failed to process PowerPoint file '{filename}'. It may be corrupted or invalid."
            ) from e

        slide_lines = []
        slide_count = len(prs.slides)
        has_visual_content = False

        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines.append(f"# Slide {idx}")

            # Process shapes in order
            for shape in slide.shapes:
                # Detect visual content (pictures, charts, diagrams, group shapes)
                if shape.shape_type in (
                    MSO_SHAPE_TYPE.PICTURE,
                    MSO_SHAPE_TYPE.CHART,
                    MSO_SHAPE_TYPE.DIAGRAM,
                    MSO_SHAPE_TYPE.GROUP,
                    MSO_SHAPE_TYPE.MEDIA,
                    MSO_SHAPE_TYPE.FREEFORM,
                    MSO_SHAPE_TYPE.IGX_GRAPHIC,
                    MSO_SHAPE_TYPE.LINKED_PICTURE,
                ) or hasattr(shape, "image"):
                    has_visual_content = True

                # Extract Text Frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_lines.append(text)

                # Extract Tables
                if shape.has_table:
                    table_md = format_markdown_table(shape.table)
                    if table_md:
                        slide_lines.append("\n### Table")
                        slide_lines.append(table_md)

            # Extract Speaker Notes
            if slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_lines.append("\n### Speaker Notes")
                            slide_lines.append(notes_text)
                except Exception as notes_err:
                    logger.debug(f"Could not read speaker notes for slide {idx}: {notes_err}")

            slide_lines.append("")  # Blank line between slides

        full_md = "\n".join(slide_lines).strip()

        warnings = []
        if has_visual_content:
            warnings.append(
                "This presentation contains visual content (pictures/charts/diagrams) that Uno cannot interpret yet."
            )

        if not full_md:
            warnings.append("No text content could be extracted from this presentation.")

        return DocumentAnalysis(
            filename=filename,
            file_type="PPTX",
            markdown=full_md,
            warnings=tuple(warnings),
            slide_count=slide_count,
        )
